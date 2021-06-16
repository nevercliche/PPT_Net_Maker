import matplotlib.pyplot as plt
import networkx as nx
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches,Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Cm



# Global Parameters
LAYOUT_KIND = 0 # 0-spring_layout  1-circular  2-random  3-shell  4-spectral
NODE_SIZE = 0.3  #Inches
LINE_SIZE = 0
TEXT_SIZE = 0
#颜色可以使用RGB
NODE_COLOR = MSO_THEME_COLOR.ACCENT_1
LINE_COLOR = MSO_THEME_COLOR.ACCENT_1
TEXT_COLOR = MSO_THEME_COLOR.ACCENT_2

def getposition(G,choice):
    if choice ==1:
        return nx.circular_layout(G)
    elif choice ==2:
        return nx.random_layout(G)
    elif choice ==3:
        return nx.shell_layout(G)
    elif choice ==4:
        return nx.spectral_layout(G)
    else :
        return nx.spring_layout(G)


#direc: 0-top 1-left 2-down 3-right (chosen by the relative position
def getdirec(node1, node2):
    direc = []
    temp_x = node2[0] - node1[0]
    temp_y = node2[1] - node1[1]
    if (temp_y >= 0 and abs(temp_x) < temp_y):
        direc.append(0)
        direc.append(2)
    elif (temp_y <= 0 and abs(temp_x) < abs(temp_y)):
        direc.append(2)
        direc.append(0)
    elif (temp_x >= 0 and abs(temp_y) < temp_x):
        direc.append(3)
        direc.append(1)
    else:
        direc.append(1)
        direc.append(3)
    return direc

def main():
    '''
    H = nx.path_graph(8)
    edges = [(0,1),(0,2),(1,3),(1,4),(2,5),(3,7),(4,7),(2,6)]
    G = nx.Graph()
    '''
    file = open('net.txt','rb')
    line = file.readline()
    edges = []#保存边关系，（序号1，序号2）

    #方便由序号找到节点名，由节点名找到对应的序号（没有找到合适的数据结构
    n = dict()#保存节点,(序号，节点名称)
    n_reverse = dict() #(节点名称，序号)

    index = 0 #记录序号
    while line:
        #edge = line.decode().split('，')
        line_de = line.decode()
        edge = re.split(r'[，\r]', line_de)
        x = edge[0]
        x_index,y_index=0,0
        if x in n_reverse:
            x_index = n_reverse.get(x)
        else:
            n[index]=x
            n_reverse[x]=index
            x_index = index
            index += 1
        y = edge[1]
        if y in n_reverse:
            y_index = n_reverse.get(y)
        else:
            n[index]=y
            n_reverse[y]=index
            y_index = index
            index += 1
        edges.append((x_index,y_index))
        line = file.readline()
    #H = nx.path_graph(len(edges))
    G = nx.Graph()
    #判断节点数量是否小于100
    if len(n) > 100:
        print('The number of nodes should be less than 100.')

    G.add_edges_from(edges)
    G.add_nodes_from(list(range(len(n))))
    pos = []


    pos=getposition(G,LAYOUT_KIND)
    #nx.draw_networkx(G)
    #plt.show()

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    nodes=[]

    '''
    for i in range(len(pos)):
        #default:  width-10 inches   height-7.46inches
        nodes.append(shapes.add_shape(MSO_SHAPE.OVAL, Inches(5+pos[i][0]*5), Inches(3.73+pos[i][1]*3.73), Inches(0.2), Inches(0.2)))
        #set the color of the cycle
        fill = nodes[i].fill
        fill.solid()
        fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        fill.fore_color.brightness = -0.25
        '''
    i = 0
    for key,value in pos.items():
        #default:  width-10 inches   height-7.46inches
        nodes.append(shapes.add_shape(MSO_SHAPE.OVAL, Inches(5+pos[key][0]*5), Inches(3.73+pos[key][1]*3.73), Inches(NODE_SIZE), Inches(NODE_SIZE)))
        #set the color of the cycle
        fill = nodes[i].fill
        fill.solid()
        fill.fore_color.theme_color = NODE_COLOR
        fill.fore_color.brightness = -0.25

        i = i+1

    for i in range(len(G.edges)-1):
        connector = shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Cm(2), Cm(2), Cm(10), Cm(10)
        )
        direc = chooseDirection(pos[edges[i][0]],pos[edges[i][1]])
        connector.begin_connect(nodes[edges[i][0]],direc[0])
        connector.end_connect(nodes[edges[i][1]],direc[1])


    prs.save('test.pptx')


if __name__ == '__main__':
    main()
